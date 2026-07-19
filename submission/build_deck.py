"""Fill the KSP Datathon 2026 submission template from the real project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

SRC = "ppt/template.pptx"
OUT = "KSP_Crime_Intelligence_Submission.pptx"

# palette
INK   = RGBColor(0x11,0x18,0x27)
NAVY  = RGBColor(0x14,0x27,0x4E)
TEAL  = RGBColor(0x0E,0x6E,0x6B)
BLUE  = RGBColor(0x21,0x4A,0xA8)
AMBER = RGBColor(0xB5,0x6C,0x08)
GREEN = RGBColor(0x1B,0x7A,0x43)
RED   = RGBColor(0xB1,0x2A,0x2A)
GREY  = RGBColor(0x5B,0x65,0x77)
CARD  = RGBColor(0xF2,0xF5,0xFA)
CARD2 = RGBColor(0xEA,0xF1,0xF0)
LINE  = RGBColor(0xD8,0xDF,0xEA)
WHITE = RGBColor(0xFF,0xFF,0xFF)
DARKUI= RGBColor(0x10,0x16,0x1F)
DARKUI2=RGBColor(0x18,0x20,0x2E)
FONT = "Calibri"
HFONT = "Calibri"

prs = Presentation(SRC)
S = prs.slides

def _set_fill(shape, color):
    if color is None:
        shape.fill.background()
    else:
        shape.fill.solid(); shape.fill.fore_color.rgb = color

def _set_line(shape, color, w=1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color; shape.line.width = Pt(w)

def box(slide, x, y, w, h, fill=CARD, line=LINE, lw=1.0, radius=0.09, shadow=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    _set_fill(sp, fill); _set_line(sp, line, lw)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    if shadow:
        el = sp._element.spPr
        # subtle outer shadow
        from pptx.oxml import parse_xml
        ns = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        shdw = parse_xml(f'<a:effectLst {ns}><a:outerShdw blurRad="90000" dist="30000" dir="5400000" rotWithShape="0"><a:srgbClr val="1A2540"><a:alpha val="18000"/></a:srgbClr></a:outerShdw></a:effectLst>')
        el.append(shdw)
    sp.text_frame.word_wrap = True
    return sp

def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts: {t, size, bold, color, italic, font, space_after, space_before, bullet, level, align}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", align)
        if p.get("space_after") is not None: para.space_after = Pt(p["space_after"])
        if p.get("space_before") is not None: para.space_before = Pt(p["space_before"])
        para.level = p.get("level", 0)
        runs = p["runs"] if "runs" in p else [p]
        for r in runs:
            run = para.add_run(); run.text = r["t"]
            f = run.font
            f.size = Pt(r.get("size", 12)); f.bold = r.get("bold", False)
            f.italic = r.get("italic", False); f.name = r.get("font", FONT)
            f.color.rgb = r.get("color", INK)
    return tb

def stat(slide, x, y, w, num, label, color=NAVY, sub=None):
    box(slide, x, y, w, 1.15, fill=CARD, line=LINE, shadow=True)
    paras=[{"runs":[{"t":num,"size":30,"bold":True,"color":color,"font":HFONT}],"align":PP_ALIGN.CENTER,"space_after":1},
           {"runs":[{"t":label,"size":10.5,"bold":True,"color":INK,"font":FONT}],"align":PP_ALIGN.CENTER,"space_after":0}]
    if sub: paras.append({"runs":[{"t":sub,"size":8.5,"color":GREY,"font":FONT}],"align":PP_ALIGN.CENTER})
    text(slide, x, y+0.12, w, 0.95, paras, anchor=MSO_ANCHOR.TOP)

def chip(slide, x, y, w, h, label, fill, tcolor=WHITE, size=9.5, bold=True):
    sp = box(slide, x, y, w, h, fill=fill, line=None, radius=0.5)
    text(slide, x, y, w, h, [{"runs":[{"t":label,"size":size,"bold":bold,"color":tcolor,"font":FONT}],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    return sp

def arrow(slide, x, y, w, h=0.0, color=NAVY):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.22) if h==0 else Inches(h))
    _set_fill(sp, color); _set_line(sp, None); sp.shadow.inherit=False
    try: sp.adjustments[0]=0.5; sp.adjustments[1]=0.55
    except Exception: pass
    return sp

def circle(slide, x, y, d, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    _set_fill(sp, fill); _set_line(sp, line); sp.shadow.inherit=False
    return sp

def line_connector(slide, x1, y1, x2, y2, color=LINE, w=1.5):
    ln = slide.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(w); ln.shadow.inherit=False
    return ln

def del_shape(sh):
    sh._element.getparent().remove(sh._element)

def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt

# =====================================================================
# SLIDE 1 — Team Details  (title bg image3: white lower third y>=3.4)
# =====================================================================
s1 = S[0]
for sh in list(s1.shapes):
    if sh.has_text_frame and "Team" in (sh.text or ""):
        del_shape(sh)
text(s1, 0.55, 3.42, 9.0, 0.6, [
    {"runs":[{"t":"KSP Crime Intelligence Platform","size":26,"bold":True,"color":NAVY,"font":HFONT}]}])
text(s1, 0.55, 3.98, 9.0, 0.35, [
    {"runs":[{"t":"Conversational AI + crime analytics over the Karnataka State Police FIR database","size":12.5,"italic":True,"color":TEAL,"font":FONT}]}])
def field(lbl, val, y, valcolor=INK):
    text(s1, 0.6, y, 8.6, 0.32, [{"runs":[
        {"t":lbl+"  ","size":12.5,"bold":True,"color":GREY,"font":FONT},
        {"t":val,"size":12.5,"bold":True,"color":valcolor,"font":FONT}]}])
field("Team name:", "________________________", 4.45)
field("Team leader name:", "________________________", 4.78)
field("Team size:", "______", 5.11)
text(s1, 4.7, 4.45, 4.85, 1.0, [
    {"runs":[{"t":"Problem Statement:  ","size":12.5,"bold":True,"color":GREY,"font":FONT},
             {"t":"PS-1","size":12.5,"bold":True,"color":AMBER,"font":FONT}]},
    {"runs":[{"t":"Conversational AI & Crime Analytics over the KSP FIR Database","size":11.5,"color":INK,"font":FONT}],"space_before":2}])
notes(s1, "Fill in team name, leader and size. Problem Statement 1: a conversational AI plus a crime-analytics engine over the Karnataka State Police FIR database.")

# =====================================================================
# Helper: subtitle under the template heading on a content slide
# =====================================================================
def subtitle(slide, txt, color=TEAL):
    text(slide, 0.55, 1.34, 9.0, 0.32, [{"runs":[{"t":txt,"size":12,"italic":True,"color":color,"font":FONT}]}])

# =====================================================================
# SLIDE 2 — Brief about the solution
# =====================================================================
s2 = S[1]
subtitle(s2, "Ask in plain language. Get only tool-verified facts, with the evidence attached.")
brief = [
 {"runs":[{"t":"Any officer — from a station SHO to an SCRB analyst — asks a question in plain "
   "language and gets an answer assembled ","size":12.5,"color":INK,"font":FONT},
   {"t":"only","size":12.5,"bold":True,"italic":True,"color":NAVY,"font":FONT},
   {"t":" from typed tool results. The assistant never writes SQL and never authors a fact; "
    "when no tool can answer, it says so instead of guessing.","size":12.5,"color":INK,"font":FONT}],"space_after":8},
 {"runs":[{"t":"Its core is an ","size":12.5,"color":INK,"font":FONT},
   {"t":"entity-resolution engine","size":12.5,"bold":True,"color":NAVY,"font":FONT},
   {"t":" that reconstructs the real person behind dozens of misspelt, transliterated and "
    "aliased FIR name variants — the person entity the source schema never stored.","size":12.5,"color":INK,"font":FONT}],"space_after":8},
 {"runs":[{"t":"On that spine sit ","size":12.5,"color":INK,"font":FONT},
   {"t":"21 tools","size":12.5,"bold":True,"color":TEAL,"font":FONT},
   {"t":" for case retrieval, co-offender networks, hotspots & trends, MO fingerprinting and "
    "compliance — each answer carrying a clickable evidence trail and enforced by role-based "
    "access control.","size":12.5,"color":INK,"font":FONT}]},
]
box(s2, 0.5, 1.75, 5.75, 3.35, fill=WHITE, line=LINE, shadow=True)
text(s2, 0.72, 1.95, 5.35, 3.0, brief, anchor=MSO_ANCHOR.TOP)
stat(s2, 6.55, 1.75, 3.0, "21", "typed tools in one chat loop", NAVY)
stat(s2, 6.55, 3.05, 3.0, "50,000", "synthetic FIRs, seed-reproducible", TEAL)
stat(s2, 6.55, 4.35, 3.0, "0", "facts authored by the LLM", AMBER, "everything is tool-sourced")
notes(s2, "The one-line pitch: grounded conversational intelligence over FIRs, with entity resolution as the differentiator and provenance on every answer.")

# =====================================================================
# SLIDE 3 — Opportunities / different / solves / USP
# =====================================================================
s3 = S[2]
subtitle(s3, "Why this is different, how it solves the problem, and what only we do.")
cols = [
 ("HOW IT'S DIFFERENT", NAVY,
  "AccusedMasterID in the KSP schema is a per-FIR row id, not a person. Every naive approach "
  "treats it as identity and silently mis-counts. We resolve the real person across FIRs — the "
  "schema's missing entity — and only then analyse."),
 ("HOW IT SOLVES IT", TEAL,
  "Officers get instant, jurisdiction-scoped answers. Repeat offenders and cross-station gangs "
  "invisible to any single station surface automatically. Chargesheet deadlines and registration "
  "delays are watched continuously."),
 ("USP", AMBER,
  "Grounded (never fabricates), Explainable (clickable provenance to the source FIR), Governed "
  "(RBAC + audit at the tool boundary, not the prompt), and Measured (entity-resolution B³ "
  "F1 = 0.687 against hidden ground truth)."),
]
cx = 0.5; cw = 2.95; gap = 0.30
for i,(h,c,body) in enumerate(cols):
    x = cx + i*(cw+gap)
    box(s3, x, 1.75, cw, 3.35, fill=CARD, line=LINE, shadow=True)
    circle(s3, x+0.28, 2.0, 0.34, c)
    text(s3, x+0.28, 2.0, 0.34, 0.34, [{"runs":[{"t":str(i+1),"size":13,"bold":True,"color":WHITE,"font":HFONT}],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s3, x+0.72, 2.02, cw-0.9, 0.4, [{"runs":[{"t":h,"size":12.5,"bold":True,"color":c,"font":HFONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s3, x+0.24, 2.55, cw-0.46, 2.4, [{"runs":[{"t":body,"size":11.5,"color":INK,"font":FONT}]}])
notes(s3, "The differentiator is entity resolution; the moat is that everything is grounded, explainable, governed and measured.")

# =====================================================================
# SLIDE 4 — Features (2 x 4 grid)
# =====================================================================
s4 = S[3]
subtitle(s4, "21 typed tools, grouped into eight capability areas.")
feats = [
 ("Conversational Q&A", "Grounded answers; refuses when no tool fits", NAVY),
 ("Entity Resolution", "The person behind dozens of name variants", TEAL),
 ("Co-offender Networks", "Gangs & cross-jurisdiction groups (Louvain)", BLUE),
 ("Hotspots & Trends", "DBSCAN hotspots, red-zone z-score, seasonality", GREEN),
 ("MO Fingerprinting", "Similar cases returned with their outcomes", AMBER),
 ("Compliance", "Chargesheet-deadline board, registration delay", RED),
 ("RBAC", "5 roles, scope enforced at the tool boundary", NAVY),
 ("Evidence & Audit", "Provenance on every answer; every call logged", TEAL),
]
gx, gy, gw, gh, hgap, vgap = 0.5, 1.72, 2.17, 1.55, 0.22, 0.22
for i,(t,d,c) in enumerate(feats):
    r,cc = divmod(i,4)
    x = gx + cc*(gw+hgap); y = gy + r*(gh+vgap)
    box(s4, x, y, gw, gh, fill=CARD, line=LINE, shadow=True)
    circle(s4, x+0.2, y+0.2, 0.22, c)
    text(s4, x+0.16, y+0.5, gw-0.3, 0.5, [{"runs":[{"t":t,"size":11.5,"bold":True,"color":c,"font":HFONT}]}])
    text(s4, x+0.16, y+0.9, gw-0.3, 0.6, [{"runs":[{"t":d,"size":9.5,"color":INK,"font":FONT}]}])
notes(s4, "Every feature is a typed tool the LLM can call; adding a tool auto-registers it into the chat.")

# =====================================================================
# SLIDE 5 — Process flow
# =====================================================================
s5 = S[4]
subtitle(s5, "The LLM only selects tools — it never writes SQL and never authors a fact.")
nodes = [
 ("Officer\nquestion", NAVY),
 ("Orchestration\nloop (LLM)", BLUE),
 ("Selects a\ntyped tool", TEAL),
 ("RBAC + k-anon\nat the boundary", AMBER),
 ("Parameterized\nSQL / ML", GREEN),
 ("Grounded answer\n+ provenance", NAVY),
]
n = len(nodes); nw=1.30; nh=1.05; y=2.1
total = 9.0; gap = (total - n*nw)/(n-1)
x = 0.5
xs=[]
for i,(lbl,c) in enumerate(nodes):
    xi = 0.5 + i*(nw+gap); xs.append(xi)
    box(s5, xi, y, nw, nh, fill=CARD, line=c, lw=1.5, shadow=True)
    text(s5, xi, y, nw, nh, [{"runs":[{"t":lbl,"size":10.5,"bold":True,"color":c,"font":FONT}],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    if i>0:
        ax = xs[i-1]+nw; arrow(s5, ax+0.02, y+nh/2-0.11, gap-0.04, color=GREY)
box(s5, 0.5, 3.75, 9.0, 1.15, fill=CARD2, line=LINE)
text(s5, 0.75, 3.9, 8.5, 0.9, [
 {"runs":[{"t":"Use case  ","size":11.5,"bold":True,"color":TEAL,"font":HFONT},
   {"t":"“Show me repeat offenders in my district and the gangs they run in.”","size":11.5,"italic":True,"color":INK,"font":FONT}],"space_after":3},
 {"runs":[{"t":"→ get_repeat_offenders → get_person_network → detect_communities, each scoped to "
   "the officer’s jurisdiction, each row traceable to its CrimeNos.","size":11,"color":GREY,"font":FONT}]}])
notes(s5, "The loop: question -> tool selection -> RBAC-checked parameterized execution -> grounded answer with provenance.")

# =====================================================================
# SLIDE 6 — Wireframe / mock (3-pane UI)
# =====================================================================
s6 = S[5]
subtitle(s6, "Dark, dense, operational — a police tool, not a consumer app.")
# app frame
box(s6, 0.5, 1.7, 9.0, 3.45, fill=DARKUI, line=RGBColor(0x2A,0x35,0x48), radius=0.03, shadow=True)
# header
box(s6, 0.5, 1.7, 9.0, 0.5, fill=DARKUI2, line=None, radius=0.03)
text(s6, 0.72, 1.72, 5.0, 0.46, [{"runs":[{"t":"KSP Crime Intelligence","size":11,"bold":True,"color":WHITE,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
chip(s6, 7.35, 1.8, 1.9, 0.3, "Acting as: SP  ▾", RGBColor(0x22,0x2C,0x3D), WHITE, 9, True)
# left chat pane
box(s6, 0.66, 2.34, 4.2, 2.68, fill=DARKUI2, line=None, radius=0.03)
chip(s6, 0.85, 2.5, 2.9, 0.34, "Who are the repeat offenders?", RGBColor(0x2B,0x4A,0x7D), WHITE, 9, False)
box(s6, 1.2, 2.95, 3.5, 0.62, fill=RGBColor(0x1A,0x22,0x31), line=None, radius=0.06)
text(s6, 1.3, 2.98, 3.3, 0.56, [{"runs":[{"t":"8 people with ≥3 FIRs in your district; top is linked to a 4-station group.","size":8.5,"color":RGBColor(0xC9,0xD2,0xE0),"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
for i,(lbl) in enumerate(["get_repeat_offenders","detect_communities"]):
    chip(s6, 1.2+i*1.75, 3.64, 1.65, 0.26, lbl, RGBColor(0x1C,0x25,0x36), RGBColor(0x8A,0x94,0xA6), 7.5, False)
box(s6, 0.85, 4.55, 3.85, 0.34, fill=RGBColor(0x16,0x1C,0x28), line=RGBColor(0x2A,0x35,0x48), radius=0.2)
text(s6, 0.95, 4.55, 3.0, 0.34, [{"runs":[{"t":"Ask the assistant…","size":8.5,"color":RGBColor(0x5B,0x65,0x77),"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
chip(s6, 4.05, 4.58, 0.55, 0.28, "Send", RGBColor(0x4C,0x8D,0xFF), RGBColor(0x06,0x12,0x24), 8.5, True)
# right pane tabs
for i,(lbl,act) in enumerate([("Evidence",False),("Network",True),("Map",False)]):
    chip(s6, 5.0+i*1.15, 2.34, 1.08, 0.3, lbl, RGBColor(0x11,0x16,0x1F) if act else RGBColor(0x18,0x20,0x2E), WHITE if act else RGBColor(0x8A,0x94,0xA6), 9, act)
box(s6, 5.0, 2.7, 4.33, 2.32, fill=RGBColor(0x0C,0x11,0x1A), line=RGBColor(0x22,0x2C,0x3D), radius=0.03)
# mini network graph
import math
cxg, cyg, rg = 7.15, 3.85, 0.72
circle(s6, cxg-0.1, cyg-0.1, 0.2, RGBColor(0x4C,0x8D,0xFF))
for k in range(6):
    a = 2*math.pi*k/6
    nx, ny = cxg+rg*math.cos(a), cyg+rg*math.sin(a)
    line_connector(s6, cxg, cyg, nx, ny, RGBColor(0x35,0x41,0x5A), 1.2)
    circle(s6, nx-0.07, ny-0.07, 0.14, RGBColor(0x2B,0x4A,0x7D), RGBColor(0x4C,0x8D,0xFF))
text(s6, 5.1, 4.72, 4.1, 0.28, [{"runs":[{"t":"co-offending network · node size = linked FIRs","size":8,"color":RGBColor(0x8A,0x94,0xA6),"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
notes(s6, "Left: streaming chat with tool chips. Right: Evidence / Network / Map panes that repaint from each answer. Role switcher drives RBAC live.")

# =====================================================================
# SLIDE 7 — Architecture
# =====================================================================
s7 = S[6]
subtitle(s7, "Everything above the database runs on Zoho Catalyst.")
layers = [
 ("FRONTEND", "React + Vite + TypeScript  —  streaming chat, evidence / network / map panes", BLUE, "Catalyst Slate"),
 ("API / ORCHESTRATION", "FastAPI tool-calling loop  —  LLM selects typed tools, never writes SQL", NAVY, "Catalyst AppSail"),
 ("TOOL LAYER (21 tools)", "RBAC + provenance + audit + k-anonymity enforced at the boundary", TEAL, "AppSail"),
 ("INTELLIGENCE", "Entity resolution · MO clustering · networks · trends  (scikit-learn, networkx)", GREEN, "AppSail / QuickML"),
 ("DATA", "PostgreSQL 16 + PostGIS + pgvector + pg_trgm  —  ksp (source) / derived (computed)", AMBER, "External managed PG"),
]
ly = 1.72; lh = 0.62; lg = 0.115; lx = 0.5; lw = 6.55
for i,(h,d,c,svc) in enumerate(layers):
    y = ly + i*(lh+lg)
    box(s7, lx, y, lw, lh, fill=CARD, line=c, lw=1.5)
    text(s7, lx+0.2, y, 2.0, lh, [{"runs":[{"t":h,"size":10,"bold":True,"color":c,"font":HFONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s7, lx+2.15, y, lw-2.3, lh, [{"runs":[{"t":d,"size":9.5,"color":INK,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
# catalyst services side column
box(s7, 7.35, 1.72, 2.15, 3.4, fill=NAVY, line=None, radius=0.05, shadow=True)
text(s7, 7.45, 1.85, 1.95, 0.3, [{"runs":[{"t":"CATALYST SERVICES","size":10,"bold":True,"color":WHITE,"font":HFONT}]}])
svcs = ["AppSail — backend","Slate — frontend","QuickML — LLM","Job Scheduling — alerts","SmartBrowz — PDF","Stratus — files","Cache · Auth"]
text(s7, 7.45, 2.25, 1.95, 2.8, [{"runs":[{"t":sv,"size":9.5,"color":RGBColor(0xCA,0xDC,0xFC),"font":FONT}],"space_after":6} for sv in svcs])
notes(s7, "Layered: Slate frontend -> AppSail FastAPI orchestration -> tool layer -> intelligence (sklearn/networkx) -> external Postgres+PostGIS+pgvector. Catalyst services wrap it; Postgres stays external because ZCQL lacks CTEs/PostGIS/pgvector.")

# =====================================================================
# SLIDE 8 — Technologies
# =====================================================================
s8 = S[7]
subtitle(s8, "Proven, open components — plus Catalyst-native AI.")
groups = [
 ("Backend", NAVY, "Python 3.12 · FastAPI · Pydantic v2 · SQLAlchemy · psycopg3"),
 ("Data", AMBER, "PostgreSQL 16 · PostGIS · pgvector · pg_trgm · fuzzystrmatch"),
 ("ML & Analytics", GREEN, "scikit-learn (HDBSCAN, DBSCAN, TF-IDF/SVD) · networkx · jellyfish · Double Metaphone · indic-transliteration"),
 ("LLM", TEAL, "Catalyst QuickML LLM Serving (Qwen 2.5) · OpenAI-compatible tool calling · prompted + native modes"),
 ("Frontend", BLUE, "React 18 · Vite · TypeScript · SVG panes (Cytoscape / MapLibre upgrade planned)"),
 ("Language & Export", RED, "Bhashini (Kannada ASR/TTS) · Catalyst SmartBrowz PDF  [planned]"),
]
gy2=1.72; gh2=1.03; vg=0.16; gw2=4.4; hg=0.2
for i,(h,c,d) in enumerate(groups):
    r,cc = divmod(i,2)
    x = 0.5 + cc*(gw2+hg); y = gy2 + r*(gh2+vg)
    box(s8, x, y, gw2, gh2, fill=CARD, line=LINE, shadow=True)
    circle(s8, x+0.22, y+0.22, 0.2, c)
    text(s8, x+0.55, y+0.15, gw2-0.7, 0.34, [{"runs":[{"t":h,"size":12,"bold":True,"color":c,"font":HFONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s8, x+0.24, y+0.52, gw2-0.46, 0.46, [{"runs":[{"t":d,"size":10,"color":INK,"font":FONT}]}])
notes(s8, "Standard, defensible stack; the only Catalyst-specific piece is the LLM via QuickML serving.")

# =====================================================================
# SLIDE 9 — Catalyst services
# =====================================================================
s9 = S[8]
subtitle(s9, "Hosting maximised on the platform partner; only the database stays external.")
svc = [
 ("AppSail", "FastAPI backend (Python runtime or Docker OCI)", NAVY),
 ("Slate", "React single-page frontend hosting", BLUE),
 ("QuickML LLM Serving", "The conversational LLM (Qwen 2.5, BYOK)", TEAL),
 ("Job Scheduling", "Proactive alert engine (deadline / red-zone)", GREEN),
 ("SmartBrowz", "Server-side PDF / briefing export", AMBER),
 ("Stratus", "Object storage for generated reports", RED),
 ("Cache", "Hot lookups & rate limiting", NAVY),
 ("Authentication", "Officer identity upstream of RBAC", TEAL),
]
gy3=1.72; gh3=0.79; vg3=0.13; gw3=4.4; hg3=0.2
for i,(h,d,c) in enumerate(svc):
    r,cc = divmod(i,2)
    x = 0.5 + cc*(gw3+hg3); y = gy3 + r*(gh3+vg3)
    box(s9, x, y, gw3, gh3, fill=CARD, line=LINE)
    box(s9, x, y, 0.14, gh3, fill=c, line=None, radius=0.0, shape=MSO_SHAPE.RECTANGLE)
    text(s9, x+0.28, y+0.09, gw3-0.4, 0.34, [{"runs":[{"t":h,"size":11.5,"bold":True,"color":c,"font":HFONT}]}])
    text(s9, x+0.28, y+0.42, gw3-0.4, 0.3, [{"runs":[{"t":d,"size":9.5,"color":INK,"font":FONT}]}])
notes(s9, "AppSail, Slate, QuickML, Job Scheduling, SmartBrowz, Stratus, Cache, Auth. Postgres stays external because ZCQL has no recursive CTEs, PostGIS or pgvector.")

# =====================================================================
# SLIDE 10 — Estimated cost (optional)
# =====================================================================
s10 = S[9]
subtitle(s10, "Indicative pilot-scale monthly estimate (one district, ~50k FIRs).")
rows = [
 ("Catalyst AppSail", "1–2 always-on instances (backend)", "Per-instance runtime"),
 ("Catalyst QuickML LLM", "Conversational queries, BYOK", "Per-token / per-call"),
 ("Catalyst Slate", "Static frontend hosting", "Low / included"),
 ("Job Scheduling", "Nightly alert sweeps", "Per-invocation"),
 ("Stratus + Cache", "Report storage + hot lookups", "Usage-based"),
 ("External PostgreSQL", "Managed PG w/ PostGIS + pgvector", "Fixed instance"),
]
box(s10, 0.5, 1.72, 9.0, 0.42, fill=NAVY, line=None)
for lbl,x,w in [("Component",0.65,3.0),("Scope",3.7,3.4),("Cost driver",7.1,2.3)]:
    text(s10, x, 1.72, w, 0.42, [{"runs":[{"t":lbl,"size":11,"bold":True,"color":WHITE,"font":HFONT}]}], anchor=MSO_ANCHOR.MIDDLE)
yy=2.14
for i,(a,b,cc) in enumerate(rows):
    fill = WHITE if i%2==0 else CARD
    box(s10, 0.5, yy, 9.0, 0.46, fill=fill, line=LINE)
    text(s10, 0.65, yy, 3.0, 0.46, [{"runs":[{"t":a,"size":10.5,"bold":True,"color":NAVY,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s10, 3.7, yy, 3.4, 0.46, [{"runs":[{"t":b,"size":10,"color":INK,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    text(s10, 7.1, yy, 2.3, 0.46, [{"runs":[{"t":cc,"size":10,"color":GREY,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.48
text(s10, 0.5, yy+0.05, 9.0, 0.4, [{"runs":[{"t":"Serverless components scale to near-zero when idle; the fixed cost is the managed "
   "Postgres instance. Figures are indicative for a pilot and finalised against Catalyst usage tiers.","size":9.5,"italic":True,"color":GREY,"font":FONT}]}])
notes(s10, "Cost is dominated by the always-on backend instance and the managed Postgres; LLM and storage are usage-based and scale with query volume.")

# =====================================================================
# SLIDE 11 — Snapshots of the prototype
# =====================================================================
s11 = S[10]
subtitle(s11, "Representative prototype screens (dark operational UI).")
# screen A: chat + evidence
def screen(x, w, title):
    box(s11, x, 1.75, w, 3.35, fill=DARKUI, line=RGBColor(0x2A,0x35,0x48), radius=0.03, shadow=True)
    box(s11, x, 1.75, w, 0.42, fill=DARKUI2, line=None, radius=0.03)
    text(s11, x+0.16, 1.75, w-0.3, 0.42, [{"runs":[{"t":title,"size":10,"bold":True,"color":WHITE,"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
screen(0.5, 4.35, "Chat  +  Evidence trail")
chip(s11, 0.72, 2.32, 3.0, 0.32, "Give me the profile of person #1864", RGBColor(0x2B,0x4A,0x7D), WHITE, 9, False)
box(s11, 0.72, 2.74, 3.9, 0.66, fill=RGBColor(0x1A,0x22,0x31), line=None, radius=0.06)
text(s11, 0.82, 2.77, 3.7, 0.6, [{"runs":[{"t":"Linked to 6 FIRs across 3 stations; earliest 2022, undetected. Confidence 0.91.","size":8.5,"color":RGBColor(0xC9,0xD2,0xE0),"font":FONT}]}], anchor=MSO_ANCHOR.MIDDLE)
text(s11, 0.72, 3.5, 3.9, 0.24, [{"runs":[{"t":"PROVENANCE — tap a CrimeNo to open the FIR","size":8,"bold":True,"color":RGBColor(0x8A,0x94,0xA6),"font":FONT}]}])
for i in range(4):
    chip(s11, 0.72+i*0.95, 3.78, 0.88, 0.3, "1002…"+str(2201+i), RGBColor(0x1C,0x25,0x36), RGBColor(0x4C,0x8D,0xFF), 7.5, False)
text(s11, 0.72, 4.2, 3.9, 0.7, [{"runs":[{"t":"get_person ✓   → 6 rows, 6 CrimeNos. RBAC: 30 out-of-scope FIRs hidden.","size":8.5,"color":RGBColor(0xA9,0xB4,0xC6),"font":FONT}]}])
# screen B: hotspots/red-zone
screen(5.15, 4.35, "Trends  +  Hotspot map")
stat_dark = [("109","MO clusters"),("13","red-zone districts"),("0.69","MO V-measure")]
for i,(nn,ll) in enumerate(stat_dark):
    bx = 5.35+i*1.38
    box(s11, bx, 2.35, 1.28, 0.9, fill=RGBColor(0x16,0x1C,0x28), line=RGBColor(0x2A,0x35,0x48), radius=0.08)
    text(s11, bx, 2.44, 1.28, 0.5, [{"runs":[{"t":nn,"size":19,"bold":True,"color":RGBColor(0x4C,0x8D,0xFF),"font":HFONT}],"align":PP_ALIGN.CENTER}])
    text(s11, bx, 2.92, 1.28, 0.3, [{"runs":[{"t":ll,"size":8,"color":RGBColor(0xA9,0xB4,0xC6),"font":FONT}],"align":PP_ALIGN.CENTER}])
box(s11, 5.35, 3.42, 3.95, 1.5, fill=RGBColor(0x0C,0x11,0x1A), line=RGBColor(0x22,0x2C,0x3D), radius=0.03)
import random as _r; _r.seed(7)
for _ in range(16):
    px = 5.5 + _r.random()*3.6; py = 3.55 + _r.random()*1.25; d = 0.08+_r.random()*0.22
    circle(s11, px, py, d, RGBColor(0xF8,0x51,0x49))
text(s11, 5.35, 4.93, 3.95, 0.2, [{"runs":[{"t":"hotspots on precise-coordinate cases only (42% coverage, stated openly)","size":7.5,"italic":True,"color":RGBColor(0x8A,0x94,0xA6),"font":FONT}]}])
notes(s11, "Left: a grounded person profile with a clickable provenance trail and RBAC hiding out-of-scope FIRs. Right: MO/trend stats and an honest hotspot map. Live screenshots can replace these once the Catalyst LLM key is wired.")

# =====================================================================
# SLIDE 12 — Performance / benchmarking
# =====================================================================
s12 = S[11]
subtitle(s12, "Measured against hidden ground truth and a 40-question evaluation.")
stat(s12, 0.5, 1.72, 2.9, "0.687", "Entity-resolution B³ F1", NAVY, "P 0.83 / R 0.59")
stat(s12, 3.55, 1.72, 2.9, "1.000", "MO cluster homogeneity", TEAL, "V-measure 0.681")
stat(s12, 6.6, 1.72, 2.9, "40/40", "orchestration eval passed", AMBER, "incl. must-refuse cases")
# detail cards
det = [
 ("Entity Resolution", NAVY, [
   "55,716 accused rows → 35,122 resolved persons",
   "+887 extra merges from the co-offending network",
   "4,395 predicted recurring vs 4,013 true",
   "4.58M candidate pairs scored in ~30 s"]),
 ("MO & System", TEAL, [
   "109 MO clusters over 50,000 FIRs (4.9% noise)",
   "21 typed tools, 158 automated tests passing",
   "Graph build 1.6 s · MO clustering ~30 s",
   "Every answer carries a provenance chain"]),
]
for i,(h,c,items) in enumerate(det):
    x = 0.5 + i*4.75
    box(s12, x, 3.1, 4.4, 2.0, fill=CARD, line=LINE, shadow=True)
    text(s12, x+0.25, 3.24, 4.0, 0.34, [{"runs":[{"t":h,"size":12.5,"bold":True,"color":c,"font":HFONT}]}])
    text(s12, x+0.25, 3.64, 4.0, 1.4, [
      {"runs":[{"t":"•  ","size":11,"bold":True,"color":c,"font":FONT},{"t":it,"size":10.5,"color":INK,"font":FONT}],"space_after":5} for it in items])
notes(s12, "Headline: ER B-cubed F1 0.687 vs hidden ground truth (the hard, measurable ML result); MO homogeneity 1.000; 40/40 grounding eval including questions that must be refused; 158 tests green.")

# =====================================================================
# SLIDE 13 — Links
# =====================================================================
s13 = S[12]
subtitle(s13, "Repository, demo and deployment.")
links = [
 ("GitHub Public Repository", "github.com/thirumurugan-12/claux", GREEN, True),
 ("Demo Video (3 minutes)", "________________________  (to be added)", BLUE, False),
 ("Deployed Link", "________________________  (Catalyst AppSail / Slate)", NAVY, False),
]
yy=2.0
for h,v,c,ok in links:
    box(s13, 0.5, yy, 9.0, 0.86, fill=CARD, line=LINE, shadow=True)
    circle(s13, 0.75, yy+0.26, 0.34, c)
    text(s13, 0.75, yy+0.26, 0.34,0.34, [{"runs":[{"t":"✓" if ok else "…","size":13,"bold":True,"color":WHITE,"font":FONT}],"align":PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
    text(s13, 1.3, yy+0.13, 7.9, 0.34, [{"runs":[{"t":h,"size":12.5,"bold":True,"color":c,"font":HFONT}]}])
    text(s13, 1.3, yy+0.47, 7.9, 0.32, [{"runs":[{"t":v,"size":11.5,"color":INK,"font":FONT}]}])
    yy += 1.02
notes(s13, "Add the 3-minute demo video and the deployed Catalyst URL before submitting.")

# =====================================================================
# SLIDE 14 — Additional / future development
# =====================================================================
s14 = S[13]
subtitle(s14, "The roadmap, and the boundaries we deliberately keep.")
box(s14, 0.5, 1.72, 5.55, 3.4, fill=CARD, line=LINE, shadow=True)
text(s14, 0.72, 1.86, 5.1, 0.34, [{"runs":[{"t":"NEXT ON THE ROADMAP","size":12,"bold":True,"color":NAVY,"font":HFONT}]}])
road = [
 "Undetected-case & offender-risk models (prioritise open cases, never predict people)",
 "Proactive early-warning alerts on Catalyst Job Scheduling",
 "Hotspot forecasting; MapLibre + Cytoscape rich panes",
 "Kannada voice via Bhashini; PDF briefings via SmartBrowz",
 "Real-FIR ingest pipeline (schema & tools are already real-data-ready)",
]
text(s14, 0.72, 2.28, 5.15, 2.7, [
  {"runs":[{"t":"•  ","size":11,"bold":True,"color":TEAL,"font":FONT},{"t":it,"size":11,"color":INK,"font":FONT}],"space_after":8} for it in road])
box(s14, 6.25, 1.72, 3.25, 3.4, fill=CARD2, line=LINE, shadow=True)
text(s14, 6.47, 1.86, 2.85, 0.34, [{"runs":[{"t":"PRINCIPLED BOUNDARIES","size":12,"bold":True,"color":RED,"font":HFONT}]}])
bnd = [
 "No individual-level crime prediction — we predict places, times and case outcomes, never people.",
 "Financial-crime integration is a declared stub: no account / txn data exists in the FIR schema, so we never fabricate it.",
 "The LLM never authors a fact.",
]
text(s14, 6.47, 2.3, 2.9, 2.7, [
  {"runs":[{"t":"•  ","size":11,"bold":True,"color":RED,"font":FONT},{"t":it,"size":10.5,"color":INK,"font":FONT}],"space_after":8} for it in bnd])
notes(s14, "Roadmap plus the ethical lines: no person-level prediction, no fabricated financial data, no LLM-authored facts.")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
